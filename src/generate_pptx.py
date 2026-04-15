from pptx import Presentation
from pptx.util import Inches, Pt
import os

def update_pptx(template_path, output_path):
    if not os.path.exists(template_path):
        print(f"Error: {template_path} not found.")
        return
    
    prs = Presentation(template_path)
    
    # --- Slide 1: Title ---
    slide = prs.slides[0]
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if "[Title of the Project]" in shape.text:
                shape.text = "Technical Debt ROI Prediction Using Survival Analysis"
            if "Name of the Student" in shape.text:
                shape.text = "Chiranjeev Rout (Regd. No. [Your ID])"
            if "Supervised By" in shape.text:
                shape.text = "Supervised By: [Your Supervisor's Name]"
            if "Group No." in shape.text:
                shape.text = "Group No.: [Your Group Number]"

    # --- Slide 3: Problem Statement ---
    slide = prs.slides[2]
    content = [
        "Modern software development faces 'Technical Debt Overload'.",
        "SonarQube generates 1000s of warnings, but which 5% are CRITICAL?",
        "Industry Gap: Current tools don't predict failure probability.",
        "Failure Cost: Average $5,600 per minute of downtime (Gartner)."
    ]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text == "":
            tf = shape.text_frame
            for line in content:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0

    # --- Slide 4: Motivations ---
    slide = prs.slides[3]
    content = [
        "Economic Compass: Bridge the gap between engineering and finance.",
        "Resource Optimization: Fix the highest-risk modules first.",
        "Business Impact: Prevent revenue loss and downtime at firms like Visa/Qualcomm.",
        "Predictive Shift: Move from 'Detect-Fix' to 'Predict-Prevent' strategy."
    ]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text == "":
            tf = shape.text_frame
            for line in content:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0

    # --- Slide 5: Literature Review ---
    slide = prs.slides[4]
    content = [
        "Lenarduzzi (2019): Established the Apache TD Dataset (33 projects).",
        "Syer (2015): Proved Survival Analysis is better for code defects than binary classification.",
        "Borg & Tornhill (2024): Quantified 'Amplified Returns' for refactoring.",
        "Digkas (2021): Identified 'Code Churn' as the #1 predictor of failure interest."
    ]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text == "":
            tf = shape.text_frame
            for line in content:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0

    # --- Slide 6: Research Gap ---
    slide = prs.slides[5]
    content = [
        "1. Duration Gap: Existing tools treat bugs as Yes/No, not 'When'.",
        "2. Economic Disconnect: No bridge between probability and dollar values.",
        "3. Static Bias: Industry tools ignore Evolutionary Churn (how code changes over time).",
        "4. AI Era: Lack of models specifically validating TD ROI in LLM-generated code."
    ]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text == "":
            tf = shape.text_frame
            for line in content:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0

    # --- Slide 7: Architecture ---
    slide = prs.slides[6]
    content = [
        "4-Layer Technical Pipeline:",
        "Layer 1: Input (Lenarduzzi SQLite DB Query)",
        "Layer 2: Feature Extraction (Churn, Complexity, Ownership)",
        "Layer 3: Modeling (Cox PH vs. Random Forest vs. Logistic Reg)",
        "Layer 4: ROI Engine (Probability x Failure Cost)"
    ]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text == "":
            tf = shape.text_frame
            for line in content:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0

    # --- Slide 9: Dataset & Tools ---
    slide = prs.slides[8]
    content = [
        "Dataset: Lenarduzzi Technical Debt Dataset (33 Projects, 78K Commits).",
        "Libraries: Lifelines (Survival), Scikit-learn (ML), Tree-sitter (Parsing).",
        "Backend: Python 3.10+, SQLite3.",
        "Frontend: Streamlit (for Interactive ROI Dashboard)."
    ]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text == "":
            tf = shape.text_frame
            for line in content:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0

    # --- Slide 10: Progress ---
    slide = prs.slides[9]
    content = [
        "Literature Survey: 100% Complete (10 papers categorized).",
        "Dataset Acquisition: 100% (td_v2.db downloaded and mapped).",
        "Project Design: 100% (4-layer pipeline architecture finalized).",
        "Initial Logic: 100% (ROI and Survival Hazard formulas derived)."
    ]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text == "":
            tf = shape.text_frame
            for line in content:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0

    # --- Slide 11: Future Work ---
    slide = prs.slides[10]
    content = [
        "Week 1-2: Feature normalization and PH Assumption testing.",
        "Week 3-4: 3-Model tournament (Cox vs. RF vs. LogReg).",
        "Week 5-6: ROI engine integration & Streamlit dashboard build.",
        "Week 7-8: Result validation and Research Manuscript submission."
    ]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text == "":
            tf = shape.text_frame
            for line in content:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0

    # --- Slide 12: References ---
    slide = prs.slides[11]
    refs = [
        "[1] Lenarduzzi et al., MSR 2019.",
        "[2] Syer et al., TSE 2015.",
        "[3] Borg & Tornhill, arXiv 2024.",
        "[4] Digkas et al., SNCS 2021.",
        "[5] Selvanayagam et al., SANER 2026."
    ]
    for shape in slide.shapes:
        if shape.has_text_frame and "References" in shape.text:
            tf = shape.text_frame
            for r in refs:
                p = tf.add_paragraph()
                p.text = r
                p.font.size = Pt(14)

    prs.save(output_path)
    print(f"PPTX saved to {output_path}")

if __name__ == "__main__":
    update_pptx("[Template] FRP Review-1 Presentation PPT  - Copy.pptx", "Progress_Review_1_Presentation_Update.pptx")
